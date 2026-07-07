# -*- coding: utf-8 -*-
"""Build web slide decks (reveal.js) from the 2025 PowerPoints.

Reads ../oldslides/*.pptx, routes every slide through classify.py,
exports deduplicated/optimized images to media/, and writes one
self-contained reveal.js page per 2026 week plus REPORT.md.

Run from slides/:  ../.venv/Scripts/python build_slides.py
"""
from __future__ import annotations

import hashlib
import html
import io
import json
import re
import sys
from pathlib import Path

import frontmatter
import yaml
from PIL import Image
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

from classify import BASE_MAP, classify

sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding="utf-8", errors="replace")

ROOT = Path(__file__).resolve().parent.parent
OLD = ROOT / "oldslides"
MEDIA = Path(__file__).resolve().parent / "media"
SITE_BASE = "https://anastasiasalter.net/HumanitiesAI"

DECKS = {  # 2025 deck number -> filename
    1: "AI_Week_One.pptx", 2: "AI_Week_Two.pptx", 3: "AI_Week_Three.pptx",
    4: "AI_Week_Four.pptx", 5: "AI_Week_Five.pptx", 6: "AI_WeekSix.pptx",
    7: "AI_Week_Seven.pptx", 8: "AI_WeekEight.pptx", 9: "AI_WeekNine.pptx",
    10: "AI_WeekTen.pptx", 11: "AIWeekEleven.pptx", 12: "AI_Week_Twelve.pptx",
    13: "AI_WeekThirteen.pptx", 14: "AI_Week_Fourteen.pptx", 15: "AI_Week_Fifteen.pptx",
}

STEMS = ["weekone", "weektwo", "weekthree", "weekfour", "weekfive", "weeksix",
         "weekseven", "weekeight", "weeknine", "weekten", "weekeleven",
         "weektwelve", "weekthirteen", "weekfourteen"]

UNIT_COLORS = {"ghosts": "#0f766e", "shells": "#b45309",
               "puppetmasters": "#365314", "coda": "#475569"}
UNIT_LABELS = {"ghosts": "Ghosts", "shells": "Shells",
               "puppetmasters": "Puppet Masters",
               "coda": "The Net Is Vast and Infinite"}

MAX_EDGE = 1500
JPEG_QUALITY = 82


def iter_pictures(shapes):
    for sh in shapes:
        if sh.shape_type == MSO_SHAPE_TYPE.GROUP:
            yield from iter_pictures(sh.shapes)
        elif sh.shape_type == MSO_SHAPE_TYPE.PICTURE:
            yield sh


def export_image(blob: bytes, cache: dict) -> str | None:
    digest = hashlib.sha1(blob).hexdigest()[:12]
    if digest in cache:
        return cache[digest]
    try:
        img = Image.open(io.BytesIO(blob))
        img.load()
    except Exception:
        return None
    if max(img.size) > MAX_EDGE:
        img.thumbnail((MAX_EDGE, MAX_EDGE), Image.LANCZOS)
    has_alpha = img.mode in ("RGBA", "LA", "P") and (
        img.mode != "P" or "transparency" in img.info)
    if has_alpha:
        name = f"{digest}.png"
        img.save(MEDIA / name, "PNG", optimize=True)
    else:
        name = f"{digest}.jpg"
        img.convert("RGB").save(MEDIA / name, "JPEG",
                                quality=JPEG_QUALITY, optimize=True)
    cache[digest] = name
    return name


def slide_texts(slide) -> tuple[str, list[tuple[int, str]]]:
    """Return (title, [(indent_level, paragraph_text), ...])."""
    title = ""
    paras: list[tuple[int, str]] = []
    for sh in slide.shapes:
        if not sh.has_text_frame:
            continue
        is_title = getattr(sh, "is_placeholder", False) and \
            sh.placeholder_format.idx == 0
        for p in sh.text_frame.paragraphs:
            text = "".join(r.text for r in p.runs).strip()
            if not text:
                continue
            if is_title and not title:
                title = text
            else:
                paras.append((p.level, text))
    return title, paras


def extract() -> tuple[list[dict], dict]:
    MEDIA.mkdir(parents=True, exist_ok=True)
    cache: dict = {}
    records = []
    for deck, fname in DECKS.items():
        prs = Presentation(str(OLD / fname))
        for idx, slide in enumerate(prs.slides):
            title, paras = slide_texts(slide)
            images = []
            for pic in iter_pictures(slide.shapes):
                try:
                    name = export_image(pic.image.blob, cache)
                except Exception:
                    name = None
                if name:
                    images.append(name)
            notes = ""
            if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
                notes = slide.notes_slide.notes_text_frame.text.strip()
            full_text = " ".join([title] + [t for _, t in paras] + [notes])
            action, target, reason = classify(deck, full_text)
            ov = OVERRIDES.get(f"{deck}:{idx}")
            if ov and action == "keep" and ov["action"] != "keep":
                action, target, reason = ov["action"], ov.get("target"), ov["reason"]
            elif ov and ov["action"] == "keep" and "case study" in ov.get("reason", ""):
                reason = ov["reason"]
            records.append({
                "deck": deck, "index": idx, "title": title, "paras": paras,
                "images": images, "notes": notes,
                "action": action, "target": target, "reason": reason,
                # the original slide 1 of every deck is its 2025 title card;
                # 2026 title slides are regenerated, so drop deck openers
                "is_deck_title": idx == 0,
            })
        print(f"deck {deck:>2} ({fname}): {len(prs.slides)} slides")
    return records, cache


def para_html(paras) -> str:
    out = []
    for level, text in paras:
        cls = f' class="lvl{level}"' if level else ""
        out.append(f"<p{cls}>{html.escape(text)}</p>")
    return "\n".join(out)


def slide_html(rec, draft: bool) -> str:
    body = para_html(rec["paras"])
    imgs = rec["images"]
    text_len = sum(len(p) for _, p in rec["paras"]) + len(rec["title"] or "")
    full = len(imgs) == 1 and text_len < 60
    parts = ['<section class="full-image">' if full else "<section>"]
    if full:
        # image-based slide: the exported image IS the slide
        if rec["title"]:
            parts.append(f"<h2 class=\"sr-only\">{html.escape(rec['title'])}</h2>")
        parts.append(f'<img class="bleed" src="media/{imgs[0]}" alt="" loading="lazy">')
        if body:
            parts.append(body)
    else:
        if rec["title"]:
            parts.append(f"<h2>{html.escape(rec['title'])}</h2>")
        img_html = ""
        if imgs:
            cls = "one" if len(imgs) == 1 else "grid"
            img_html = f'<div class="imgs {cls}">' + "".join(
                f'<img src="media/{n}" alt="" loading="lazy">' for n in imgs) + "</div>"
        if body and img_html:
            parts.append(f'<div class="cols"><div class="txt">{body}</div>{img_html}</div>')
        else:
            parts.append(body or img_html)
    if "case study" in rec["reason"]:
        parts.append('<p class="hist-note">Sora was discontinued in 2026 — '
                     'retained here as this week\'s hype-cycle case study.</p>')
    if draft and (rec["action"] == "move" or "case study" in rec["reason"]):
        parts.append(f'<div class="prov">from the 2025 Week {rec["deck"]} deck '
                     f'— {html.escape(rec["reason"])}</div>')
    if rec["notes"]:
        parts.append(f'<aside class="notes">{html.escape(rec["notes"])}</aside>')
    parts.append("</section>")
    return "\n".join(parts)


def title_slide(meta: dict, week_num: int) -> str:
    unit = meta["unit"]
    label = UNIT_LABELS[unit]
    color = UNIT_COLORS[unit]
    module = meta["module"]
    dates = f"{meta['week_start'].strftime('%B %-d').replace(' 0', ' ')}" \
        if False else meta["week_start"].strftime("%B %d, %Y")
    return f"""<section class="title-slide">
<span class="badge" style="background:{color}">{label}</span>
<h1>{html.escape(module)}</h1>
<p class="dossier">[ ENG 6806 // FALL 2026 // WEEK OF {dates.upper()} ]</p>
</section>"""


def section_slide(text: str) -> str:
    return f'<section class="section-slide"><h2>{html.escape(text)}</h2></section>'


def deck_page(week_num: int, stem: str, meta: dict, sections: list[str],
              draft: bool) -> str:
    ribbon = ('<div class="draft-ribbon">DRAFT — under review</div>'
              if draft else "")
    return f"""<!DOCTYPE html>
<html lang="en">
<head>
<meta charset="utf-8">
<meta name="viewport" content="width=device-width, initial-scale=1">
<title>{html.escape(meta['module'])} — Slides</title>
<link rel="stylesheet" href="../assets/reveal/dist/reset.css">
<link rel="stylesheet" href="../assets/reveal/dist/reveal.css">
<link rel="stylesheet" href="gits-slides.css">
<link rel="preconnect" href="https://fonts.googleapis.com">
<link href="https://fonts.googleapis.com/css2?family=Rajdhani:wght@500;700&family=IBM+Plex+Sans:ital,wght@0,400;0,600;1,400&family=IBM+Plex+Mono:wght@400;600&display=swap" rel="stylesheet">
</head>
<body>
{ribbon}
<div class="reveal"><div class="slides">
{chr(10).join(sections)}
</div></div>
<script src="../assets/reveal/dist/reveal.js"></script>
<script src="../assets/reveal/plugin/notes/notes.js"></script>
<script>
Reveal.initialize({{ hash: true, slideNumber: 'c/t',
  plugins: [ RevealNotes ] }});
</script>
</body>
</html>"""


OVERRIDES: dict = {}
PATCHES: dict = {}


def main() -> int:
    global OVERRIDES, PATCHES
    here = Path(__file__).resolve().parent
    ov_path = here / "overrides.json"
    if ov_path.exists():
        OVERRIDES = json.loads(ov_path.read_text(encoding="utf-8"))
        print(f"loaded {len(OVERRIDES)} vision overrides")
    p_path = here / "patches.json"
    if p_path.exists():
        PATCHES = json.loads(p_path.read_text(encoding="utf-8"))
        print(f"loaded patches for {len(PATCHES)} decks")
    decks_yml = here / "decks.yml"
    if decks_yml.exists():
        flags = yaml.safe_load(decks_yml.read_text(encoding="utf-8"))
    else:
        flags = {s: True for s in STEMS}
        decks_yml.write_text(yaml.safe_dump(flags, sort_keys=False),
                             encoding="utf-8")

    records, cache = extract()
    (here / "manifest.json").write_text(
        json.dumps(records, indent=1, default=str), encoding="utf-8")

    report = ["# Slide conversion report (2025 → 2026)", "",
              "Every cut and relocation below is machine-made and awaiting "
              "instructor review. Draft ribbons stay on until each deck's "
              "flag in `slides/decks.yml` is set to `false` and the decks "
              "are regenerated.", ""]

    for week_num, stem in enumerate(STEMS, start=1):
        meta = frontmatter.load(str(ROOT / f"{stem}.md")).metadata["canvas"]
        draft = bool(flags.get(stem, True))
        native = [r for r in records
                  if r["action"] == "keep" and r["target"] == week_num
                  and not r["is_deck_title"]]
        moved = [r for r in records
                 if r["action"] == "move" and r["target"] == week_num
                 and not r["is_deck_title"]]
        cuts = [r for r in records
                if r["action"] == "cut" and BASE_MAP[r["deck"]] == week_num
                and not r["is_deck_title"]]

        sections = [title_slide(meta, week_num)]
        if week_num == 10:
            spine = [r for r in native if r["deck"] == 11]
            stretch = [r for r in native if r["deck"] == 10]
            sections += [slide_html(r, draft) for r in spine]
            if stretch:
                sections.append(section_slide(
                    "Optional Stretch: Rebuilding as an AI App"))
                sections += [slide_html(r, draft) for r in stretch]
        else:
            sections += [slide_html(r, draft) for r in native]
        if moved:
            sections.append(section_slide("Readings: moved into this week for 2026"))
            sections += [slide_html(r, draft) for r in moved]

        for op in PATCHES.get(stem, []):
            if op["op"] == "replace":
                sections[op["at"] - 1] = op["html"]
            elif op["op"] == "replace_range":
                sections[op["from"] - 1:op["to"]] = op["html"]
            elif op["op"] == "drop":
                del sections[op["at"] - 1]

        (here / f"{stem}.html").write_text(
            deck_page(week_num, stem, meta, sections, draft), encoding="utf-8")

        report.append(f"## Week {week_num} — {meta['module']}")
        report.append(f"- slides: {len(sections)} total "
                      f"({len(native)} native, {len(moved)} moved in, "
                      f"{len(cuts)} cut from source deck(s), "
                      f"draft={'yes' if draft else 'no'})")
        for r in moved:
            report.append(f"  - MOVED from 2025 deck {r['deck']} slide "
                          f"{r['index']+1}: \"{(r['title'] or '(untitled)')[:60]}\" "
                          f"— {r['reason']}")
        for r in cuts:
            report.append(f"  - CUT (2025 deck {r['deck']} slide {r['index']+1}): "
                          f"\"{(r['title'] or '(untitled)')[:60]}\" — {r['reason']}")
        report.append("")

    orphan_cuts = [r for r in records if r["action"] == "cut"]
    report.append(f"Total slides cut: {len(orphan_cuts)}; "
                  f"total images exported: {len(cache)}")
    (here / "REPORT.md").write_text("\n".join(report), encoding="utf-8")
    print("wrote", len(STEMS), "decks; report at slides/REPORT.md")
    return 0


if __name__ == "__main__":
    sys.exit(main())
